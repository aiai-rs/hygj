import os
import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image, ImageDraw
import matplotlib
matplotlib.use('Agg')  # 无头模式，适合Render
import matplotlib.pyplot as plt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import io
import tempfile

def excel_to_image(file_path, output_dir):
    """xlsx/xls/xlsm/csv -> PNG (每个sheet一张图)"""
    os.makedirs(output_dir, exist_ok=True)
    try:
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
            img_buffer = io.BytesIO()
            fig, ax = plt.subplots(figsize=(12, 8))
            ax.axis('tight')
            ax.axis('off')
            table = ax.table(cellText=df.values, colLabels=df.columns, cellLoc='center', loc='center')
            table.auto_set_font_size(False)
            table.set_fontsize(10)
            table.scale(1, 2)
            plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=150)
            img_buffer.seek(0)
            img = Image.open(img_buffer)
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            img.save(os.path.join(output_dir, f"{base_name}.png"))
            plt.close(fig)
        else:
            wb = load_workbook(file_path, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                data = [[cell.value if cell.value is not None else '' for cell in row] for row in ws.iter_rows(values_only=True)]
                if data:
                    img_buffer = io.BytesIO()
                    fig, ax = plt.subplots(figsize=(12, 8))
                    ax.axis('tight')
                    ax.axis('off')
                    if len(data) > 1:
                        table = ax.table(cellText=data[1:], colLabels=data[0], cellLoc='center', loc='center')
                        table.auto_set_font_size(False)
                        table.set_fontsize(8)
                        table.scale(1, 1.5)
                    plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=150)
                    img_buffer.seek(0)
                    img = Image.open(img_buffer)
                    img.save(os.path.join(output_dir, f"{sheet_name}.png"))
                    plt.close(fig)
    except Exception as e:
        print(f"Excel转换错误: {e}")
        return None
    return output_dir

def ppt_to_image(file_path, output_dir):
    """pptx/pot -> PNG (每页一张图，简单文本提取)"""
    os.makedirs(output_dir, exist_ok=True)
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            # 提取文本
            text_content = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_content += shape.text_frame.text + "\n"
            
            # 创建简单图片
            img = Image.new('RGB', (800, 600), color='white')
            draw = ImageDraw.Draw(img)
            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)  # Render默认字体
            except:
                font = ImageFont.load_default()
            draw.text((10, 10), f"Slide {i+1}: {text_content[:200]}...", fill='black', font=font)
            img.save(os.path.join(output_dir, f"slide_{i+1}.png"))
    except Exception as e:
        print(f"PPT转换错误: {e}")
        return None
    return output_dir

def swf_to_image(file_path, output_dir):
    """swf -> PNG (Render free tier不支持，临时跳过)"""
    # os.makedirs(output_dir, exist_ok=True)
    # try:
    #     base = os.path.splitext(os.path.basename(file_path))[0]
    #     cmd = f"pdf2png -s {output_dir}/{base} {file_path}"
    #     os.system(cmd)  # 输出多帧PNG
    #     return output_dir
    # except Exception as e:
    #     print(f"SWF转换错误: {e}")
    #     return None
    print("SWF 支持暂不可用")  # 临时日志
    return None  # 跳过，返回空

def convert_file(file_path):
    """主转换函数：返回图片路径列表"""
    ext = os.path.splitext(file_path)[1].lower()
    images = []
    with tempfile.TemporaryDirectory() as temp_output:
        if ext in ['.xlsx', '.xls', '.xlsm', '.csv']:
            output_dir = excel_to_image(file_path, temp_output)
        elif ext in ['.pptx', '.pot']:
            output_dir = ppt_to_image(file_path, temp_output)
        elif ext == '.swf':
            output_dir = swf_to_image(file_path, temp_output)
        else:
            return []  # 不支持
        
        if output_dir and os.path.exists(output_dir):
            images = [os.path.join(output_dir, f) for f in os.listdir(output_dir) if f.endswith('.png')]
    return images
